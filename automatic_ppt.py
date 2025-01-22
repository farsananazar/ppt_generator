# import os

# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")
    

# #for pdf

# from PyPDF2 import PdfReader

# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text() 
#     return text

# #for text

# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()
    

# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import word_tokenize

# nltk.download('punkt')
# nltk.download('stopwords')

# def preprocess_text(text):
#     tokens = word_tokenize(text)
#     stop_words = set(stopwords.words('english'))
#     clean_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#     return " ".join(clean_tokens)


# #text segmentation

# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans

# def segment_text(text):
#     model = SentenceTransformer('all-mpnet-base-v2')
#     sentences = text.split('.')
#     embeddings = model.encode(sentences)
#     num_clusters = 5  # Adjust as needed   
#     clustering_model = KMeans(n_clusters=num_clusters)
#     clustering_model.fit(embeddings)
#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
#     return clustered_sentences

# # summerization and key point generation

# from transformers import pipeline

# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def generate_summary(text):
#     return summarizer(text, max_length=50, min_length=25, do_sample=False)[0]['summary_text']


# #title generation

# import openai

# def generate_title(text):
#     openai.api_key = "YOUR_API_KEY"
#     response = openai.Completion.create(
#         engine="text-davinci-003",
#         prompt=f"Generate a title for: {text}",
#         max_tokens=10
#     )
#     return response['choices'][0]['text'].strip()





# #slide creation

# from pptx import Presentation

# def create_ppt(slides_data):
#     prs = Presentation()
#     for slide_data in slides_data:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = slide_data['title']
#         content.text = "\n".join(slide_data['key_points'])
#     prs.save('presentation.pptx')



# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize, word_tokenize
# from pptx import Presentation

# nltk.download('punkt')
# nltk.download('stopwords')

# # Step 1: Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Step 2: Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text() 
#     return text

# # Step 3: Preprocess text
# def preprocess_text(text):
#     tokens = word_tokenize(text)
#     stop_words = set(stopwords.words('english'))
#     clean_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#     return " ".join(clean_tokens)

# # Step 4: Segment text into points
# def extract_points(text):
#     sentences = sent_tokenize(text)
#     points = []
#     for sentence in sentences:
#         if len(sentence.split()) > 5:  # Only include sentences with meaningful length
#             points.append(sentence.strip())
#     return points

# # Step 5: Create PowerPoint with slides containing points
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     else:
#         with open(file_path, 'r', encoding='utf-8') as file:
#             text = file.read()
    
#     # Preprocess and extract points
#     clean_text = preprocess_text(text)
#     points = extract_points(clean_text)

#     # Initialize PowerPoint
#     prs = Presentation()
#     slide = None
#     max_points_per_slide = 5

#     # Add slides with points
#     for i, point in enumerate(points):
#         if i % max_points_per_slide == 0:  # Start a new slide after max points
#             slide = prs.slides.add_slide(prs.slide_layouts[1])
#             title = slide.shapes.title
#             title.text = f"Slide {i // max_points_per_slide + 1}"
#             content = slide.placeholders[1]
#             content.text = ""
#         content.text += f"- {point}\n"

#     # Save PowerPoint
#     ppt_file = "presentation_with_points.pptx"
#     prs.save(ppt_file)
#     print(f"Presentation saved as {ppt_file}")

# # Example usage
# file_path = "globel_warming.pdf"         # Replace with your PDF file path
# create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import word_tokenize, sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     tokens = word_tokenize(text)
#     stop_words = set(stopwords.words('english'))
#     clean_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#     return " ".join(clean_tokens)

# # Segment text dynamically
# def segment_text(text, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
#     sentences = sent_tokenize(text)
    
#     # Ensure there are enough sentences for clustering
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters to the number of sentences
    
#     if len(sentences) == 0:
#         raise ValueError("No valid sentences found in the text for segmentation.")
    
#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters)
#     clustering_model.fit(embeddings)
    
#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Summarize text
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def generate_summary(text):
#     return summarizer(text, max_length=50, min_length=25, do_sample=False)[0]['summary_text']

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Preprocess and segment text
#     clean_text = preprocess_text(text)
#     clusters = segment_text(clean_text, num_clusters=5)
    
#     # Generate slides data
#     slides_data = []
#     for cluster, sentences in clusters.items():
#         slide_data = {
#             'title': f"Cluster {cluster + 1}",
#             'key_points': sentences
#         }
#         slides_data.append(slide_data)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for slide_data in slides_data:
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = slide_data['title']
#         content.text = "\n".join(slide_data['key_points'])
    
#     # Save the presentation
#     output_file = 'presentation_with_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with the actual file name in the same directory
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import word_tokenize, sent_tokenize
# from sentence_transformers import SentenceTransformer
# from transformers import pipeline
# from pptx import Presentation

# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     tokens = word_tokenize(text)
#     stop_words = set(stopwords.words('english'))
#     clean_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#     return " ".join(clean_tokens)

# # Summarize long paragraphs into concise points
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def summarize_paragraph(paragraph):
#     if len(paragraph.split()) > 50:  # Only summarize if the paragraph is long
#         summary = summarizer(paragraph, max_length=50, min_length=20, do_sample=False)
#         return summary[0]['summary_text']
#     return paragraph

# # Segment text into paragraphs and summarize
# def segment_and_summarize_text(text):
#     paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
#     summarized_paragraphs = [summarize_paragraph(p) for p in paragraphs]
#     return summarized_paragraphs

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Preprocess and segment text
#     clean_text = preprocess_text(text)
#     summarized_paragraphs = segment_and_summarize_text(clean_text)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for idx, paragraph in enumerate(summarized_paragraphs):
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = f"Slide {idx + 1}"
#         content.text = paragraph
    
#     # Save the presentation
#     output_file = 'presentation_with_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with the actual file name in the same directory
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from pptx import Presentation

# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Segment text dynamically into clusters
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Preprocess text and segment it
#     sentences = preprocess_text(text)
#     clusters = segment_text(sentences, num_clusters=5)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = f"Slide {cluster + 1}"  # Title for each cluster
#         content.text = "\n".join(sentences)  # Key points for the slide
    
#     # Save the presentation
#     output_file = 'presentation_with_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)



# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# # Download NLTK data
# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Summarize sentences for meaningful points
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def summarize_sentences(sentences):
#     summarized_points = []
#     for sentence in sentences:
#         if len(sentence.split()) > 20:  # Summarize if sentence is long
#             summary = summarizer(sentence, max_length=30, min_length=10, do_sample=False)
#             summarized_points.append(summary[0]['summary_text'])
#         else:
#             summarized_points.append(sentence)
#     return summarized_points

# # Group sentences into clusters based on topics
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Generate headings for clusters
# def generate_headings(clusters):
#     headings = []
#     for cluster_id, sentences in clusters.items():
#         combined_text = " ".join(sentences)
#         if len(combined_text.split()) > 20:
#             summary = summarizer(combined_text, max_length=10, min_length=5, do_sample=False)
#             headings.append(summary[0]['summary_text'])
#         else:
#             headings.append(f"Topic {cluster_id + 1}")
#     return headings

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Preprocess text
#     sentences = preprocess_text(text)
    
#     # Summarize sentences for meaningful points
#     summarized_sentences = summarize_sentences(sentences)
    
#     # Segment text into clusters
#     clusters = segment_text(summarized_sentences, num_clusters=5)
    
#     # Generate headings for clusters
#     headings = generate_headings(clusters)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster_id, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = headings[cluster_id]  # Use generated heading as the slide title
#         content.text = "\n".join(sentences)  # Add summarized points to the slide content
    
#     # Save the presentation
#     output_file = 'presentation_with_meaningful_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# # Download NLTK data
# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Summarize sentences for meaningful points
# summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# def summarize_paragraphs(paragraphs):
#     summarized_points = []
#     for paragraph in paragraphs:
#         if len(paragraph.split()) > 50:  # Summarize if paragraph is long
#             summary = summarizer(paragraph, max_length=40, min_length=15, do_sample=False)
#             summarized_points.append(summary[0]['summary_text'])
#         else:
#             summarized_points.append(paragraph)
#     return summarized_points

# # Group sentences into clusters based on topics
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Generate headings for clusters
# def generate_headings(clusters):
#     headings = []
#     for cluster_id, sentences in clusters.items():
#         combined_text = " ".join(sentences)
#         if len(combined_text.split()) > 20:
#             summary = summarizer(combined_text, max_length=10, min_length=5, do_sample=False)
#             headings.append(summary[0]['summary_text'])
#         else:
#             headings.append(f"Topic {cluster_id + 1}")
#     return headings

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Split text into paragraphs
#     paragraphs = text.split('\n\n')
    
#     # Summarize paragraphs into key points
#     summarized_paragraphs = summarize_paragraphs(paragraphs)
    
#     # Segment text into clusters
#     clusters = segment_text(summarized_paragraphs, num_clusters=5)
    
#     # Generate headings for clusters
#     headings = generate_headings(clusters)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster_id, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = headings[cluster_id]  # Use generated heading as the slide title
#         content.text = "\n".join(sentences)  # Add summarized points to the slide content
    
#     # Save the presentation
#     output_file = 'presentation_with_meaningful_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# # Download NLTK data
# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Summarize sentences for meaningful points
# summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# def summarize_paragraphs(paragraphs):
#     summarized_points = []
#     for paragraph in paragraphs:
#         if len(paragraph.split()) > 50:  # Summarize if paragraph is long
#             summary = summarizer(paragraph, max_length=50, min_length=20, do_sample=False)
#             summarized_points.append(summary[0]['summary_text'])
#         else:
#             summarized_points.append(paragraph)
#     return summarized_points

# # Group sentences into clusters based on topics
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Generate headings for clusters
# def generate_headings(clusters):
#     headings = []
#     for cluster_id, sentences in clusters.items():
#         combined_text = " ".join(sentences)
#         if len(combined_text.split()) > 20:
#             summary = summarizer(combined_text, max_length=10, min_length=5, do_sample=False)
#             headings.append(summary[0]['summary_text'])
#         else:
#             headings.append(f"Topic {cluster_id + 1}")
#     return headings

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
    
#     # Split text into paragraphs
#     paragraphs = text.split('\n\n')
    
#     # Summarize paragraphs into key points
#     summarized_paragraphs = summarize_paragraphs(paragraphs)
    
#     # Segment text into clusters
#     clusters = segment_text(summarized_paragraphs, num_clusters=5)
    
#     # Generate headings for clusters
#     headings = generate_headings(clusters)
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster_id, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = headings[cluster_id]  # Use generated heading as the slide title
#         content.text = "\n".join([f"- {sentence}" for sentence in sentences])  # Add bullet points
    
#     # Save the presentation
#     output_file = 'presentation_with_meaningful_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\\Downloads\\globel_warming.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# # Download NLTK data
# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract titles from PDF
# def extract_titles_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     titles = []
#     for page in reader.pages:
#         if '/Title' in page.keys():
#             titles.append(page['/Title'])
#     return titles

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Summarize sentences for meaningful points
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def summarize_sentences(sentences):
#     summarized_points = []
#     for sentence in sentences:
#         if len(sentence.split()) > 20:  # Summarize if sentence is long
#             summary = summarizer(sentence, max_length=30, min_length=10, do_sample=False)
#             summarized_points.append(summary[0]['summary_text'])
#         else:
#             summarized_points.append(sentence)
#     return summarized_points

# # Group sentences into clusters based on topics
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42, n_init=10)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Generate headings for clusters
# def generate_headings(clusters):
#     headings = []
#     for cluster_id, sentences in clusters.items():
#         combined_text = " ".join(sentences)
#         if len(combined_text.split()) > 20:
#             summary = summarizer(combined_text, max_length=10, min_length=5, do_sample=False)
#             headings.append(summary[0]['summary_text'])
#         else:
#             headings.append(f"Topic {cluster_id + 1}")
#     return headings

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#         titles = extract_titles_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
#         titles = None
    
#     # Preprocess text
#     sentences = preprocess_text(text)
    
#     # Summarize sentences for meaningful points
#     summarized_sentences = summarize_sentences(sentences)
    
#     # Segment text into clusters
#     clusters = segment_text(summarized_sentences, num_clusters=min(5, len(summarized_sentences)))
    
#     # Generate headings for clusters (use extracted titles for PDFs if available)
#     if file_type == 'pdf' and titles:
#         headings = titles[:len(clusters)]  # Use extracted titles if they exist
#     else:
#         headings = generate_headings(clusters)
    
#     # Ensure headings match the number of clusters
#     if len(headings) < len(clusters):
#         headings.extend([f"Topic {i+1}" for i in range(len(headings), len(clusters))])
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster_id, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = headings[cluster_id]  # Use generated or extracted heading as the slide title
#         content.text = "\n".join(f"- {sentence}" for sentence in sentences)  # Add bullet points to the slide content
    
#     # Save the presentation
#     output_file = 'presentation_with_meaningful_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\FARSNA NAZAR RESUME.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)


# import os
# from PyPDF2 import PdfReader
# import nltk
# from nltk.corpus import stopwords
# from nltk.tokenize import sent_tokenize
# from sentence_transformers import SentenceTransformer
# from sklearn.cluster import KMeans
# from transformers import pipeline
# from pptx import Presentation

# # Download NLTK data
# nltk.download('punkt')
# nltk.download('stopwords')

# # Validate file type
# def validate_file(file_path):
#     if file_path.endswith('.pdf'):
#         return 'pdf'
#     elif file_path.endswith('.txt'):
#         return 'text'
#     else:
#         raise ValueError("Unsupported file format. Only PDF and TXT are allowed.")

# # Extract text from PDF
# def extract_text_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     text = ""
#     for page in reader.pages:
#         text += page.extract_text()
#     return text

# # Extract titles from PDF
# def extract_titles_from_pdf(file_path):
#     reader = PdfReader(file_path)
#     titles = []
#     for page in reader.pages:
#         if '/Title' in page.keys():
#             titles.append(page['/Title'])
#     return titles

# # Extract text from TXT
# def extract_text_from_txt(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         return file.read()

# # Preprocess text
# def preprocess_text(text):
#     stop_words = set(stopwords.words('english'))
#     sentences = sent_tokenize(text)
#     clean_sentences = []
#     for sentence in sentences:
#         tokens = sentence.split()
#         filtered_tokens = [token for token in tokens if token.isalnum() and token.lower() not in stop_words]
#         clean_sentences.append(" ".join(filtered_tokens))
#     return clean_sentences

# # Summarize sentences for meaningful points
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def summarize_sentences(sentences):
#     summarized_points = []
#     for sentence in sentences:
#         if len(sentence.split()) > 20:  # Summarize if sentence is long
#             summary = summarizer(sentence, max_length=30, min_length=10, do_sample=False)
#             summarized_points.append(summary[0]['summary_text'])
#         else:
#             summarized_points.append(sentence)
#     return summarized_points

# # Group sentences into clusters based on topics
# def segment_text(sentences, num_clusters=5):
#     model = SentenceTransformer('all-mpnet-base-v2')
    
#     if len(sentences) < num_clusters:
#         num_clusters = len(sentences)  # Adjust clusters if fewer sentences exist

#     embeddings = model.encode(sentences)
#     clustering_model = KMeans(n_clusters=num_clusters, random_state=42, n_init=10)
#     clustering_model.fit(embeddings)

#     clustered_sentences = {}
#     for idx, label in enumerate(clustering_model.labels_):
#         clustered_sentences.setdefault(label, []).append(sentences[idx])
    
#     return clustered_sentences

# # Generate headings for clusters
# def generate_headings(clusters):
#     headings = []
#     for cluster_id, sentences in clusters.items():
#         combined_text = " ".join(sentences)
#         if len(combined_text.split()) > 20:
#             summary = summarizer(combined_text, max_length=10, min_length=5, do_sample=False)
#             headings.append(summary[0]['summary_text'])
#         else:
#             headings.append(f"Topic {cluster_id + 1}")
#     return headings

# # Create PowerPoint presentation
# def create_ppt_with_points(file_path):
#     file_type = validate_file(file_path)
    
#     if file_type == 'pdf':
#         text = extract_text_from_pdf(file_path)
#         titles = extract_titles_from_pdf(file_path)
#     elif file_type == 'text':
#         text = extract_text_from_txt(file_path)
#         titles = None
    
#     # Preprocess text
#     sentences = preprocess_text(text)
    
#     # Summarize sentences for meaningful points
#     summarized_sentences = summarize_sentences(sentences)
    
#     # Segment text into clusters
#     clusters = segment_text(summarized_sentences, num_clusters=min(5, len(summarized_sentences)))
    
#     # Generate headings for clusters (use extracted titles for PDFs if available)
#     if file_type == 'pdf' and titles:
#         headings = titles[:len(clusters)]  # Use extracted titles if they exist
#     else:
#         headings = generate_headings(clusters)
    
#     # Ensure headings match the number of clusters
#     if len(headings) < len(clusters):
#         headings.extend([f"Topic {i+1}" for i in range(len(headings), len(clusters))])
    
#     # Create PowerPoint slides
#     prs = Presentation()
#     for cluster_id, sentences in clusters.items():
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         title = slide.shapes.title
#         content = slide.placeholders[1]
#         title.text = headings[cluster_id]  # Use generated or extracted heading as the slide title
#         content.text = "\n".join(f"- {sentence}" for sentence in sentences)  # Add bullet points to the slide content
    
#     # Save the presentation
#     output_file = 'presentation_with_meaningful_points.pptx'
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# # Main execution
# if __name__ == "__main__":
#     # Add the path to your PDF or TXT file here
#     file_path = "D:\Downloads\globel_warming.pdf"  # Replace with your actual file name
#     create_ppt_with_points(file_path)

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from gensim.summarization import summarize

def extract_text_from_pdf(pdf_path):
    text = ""
    pdf_document = fitz.open(pdf_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text

def extract_titles_from_pdf(pdf_path):
    titles = []
    pdf_document = fitz.open(pdf_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        blocks = page.get_text("blocks")
        for block in blocks:
            if block[4] == 0:
                titles.append(block[4])
    return titles

def segment_text(text):
    # Segment text into topics
    # This is a simple example, you might want to use more advanced techniques
    segments = text.split("\n\n")
    return segments

def generate_key_points(segment):
    summary = summarize(segment, word_count=50)
    key_points = summary.split("\n")
    return key_points

def generate_presentation_from_text(text_path):
    with open(text_path, "r") as file:
        text = file.read()
    
    segments = segment_text(text)
    prs = Presentation()

    for segment in segments:
        key_points = generate_key_points(segment)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Title for Segment"  # Generate appropriate title
        content.text = "\n".join(key_points)

    prs.save("generated_presentation.pptx")

def generate_presentation_from_pdf(pdf_path):
    text = extract_text_from_pdf(pdf_path)
    titles = extract_titles_from_pdf(pdf_path)
    segments = segment_text(text)
    prs = Presentation()

    for i, segment in enumerate(segments):
        key_points = generate_key_points(segment)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = titles[i] if i < len(titles) else "Title for Segment"
        content.text = "\n".join(key_points)

    prs.save("generated_presentation.pptx")

# Example usage
# Replace 'example.txt' and 'example.pdf' with the path to your text and PDF files
# text_file_path = "path/to/your/example.txt"
pdf_file_path = "D:\Downloads\globel_warming.pdf"

# Generate presentation from text file
# generate_presentation_from_text(text_file_path)

# Generate presentation from PDF file
generate_presentation_from_pdf(pdf_file_path)


