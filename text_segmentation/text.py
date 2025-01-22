# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# import matplotlib.pyplot as plt
# from summa.summarizer import summarize  # Importing summa's summarizer

# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix

# def summarize_text(segmented_text, ratio=0.3):
#     """
#     Summarize each segment using the TextRank algorithm.
#     """
#     summarized_segments = []
#     for segment in segmented_text:
#         # Summarize the text using summa's summarize() function
#         summary = summarize(segment, ratio=ratio)
#         summarized_segments.append(summary if summary else segment)  # Fallback to original text if summary fails
#     return summarized_segments

# # Input: List of paragraphs
# input_paragraphs = [
#     "Machine learning has revolutionized industries like healthcare, finance, and retail.",
#     "Climate change is one of the most pressing challenges of our time.",
#     "Efforts to reduce carbon emissions and adopt renewable energy are critical.",
#     "Artificial intelligence is increasingly being used to solve complex global problems."
# ]

# # Perform text segmentation
# segmented_text, similarity_matrix = segment_text(input_paragraphs, threshold=0.5)

# # Summarize each segment
# summarized_segments = summarize_text(segmented_text, ratio=0.5)

# # Print summarized text
# print("Summarized Segments:")
# for i, summary in enumerate(summarized_segments, 1):
#     print(f"Summary {i}: {summary}")

#$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$

# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# import matplotlib.pyplot as plt
# from summa.summarizer import summarize  # Importing summa's summarizer
# import PyPDF2  # For reading PDF files

# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix

# def summarize_text(segmented_text, ratio=0.3):
#     """
#     Summarize each segment using the TextRank algorithm.
#     """
#     summarized_segments = []
#     for segment in segmented_text:
#         # Summarize the text using summa's summarize() function
#         summary = summarize(segment, ratio=ratio)
#         summarized_segments.append(summary if summary else segment)  # Fallback to original text if summary fails
#     return summarized_segments

# def main(file_path, threshold=0.5, ratio=0.3):
#     """
#     Main function to read file, segment text, and summarize.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Summarize each segment
#     summarized_segments = summarize_text(segmented_text, ratio=ratio)

#     # Print summarized text
#     print("Summarized Segments:")
#     for i, summary in enumerate(summarized_segments, 1):
#         print(f"Summary {i}: {summary}")

# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "globel_warming.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path)

#$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$

# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from summa.summarizer import summarize
# import PyPDF2


# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix

# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         # Split the segment into sentences or logical points
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def main(file_path, threshold=0.5):
#     """
#     Main function to read file, segment text, and convert to bullet points.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Print bullet points for each segment
#     print("Bullet Points for Segmented Text:")
#     for i, bullets in enumerate(bullet_points_segments, 1):
#         print(f"\nSegment {i}:")
#         for bullet in bullets:
#             print(bullet)

# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "globel_warming.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path)
# ****************************************************************************************************************************8

# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2


# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix

# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         # Split the segment into sentences or logical points
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def create_presentation(bullet_points_segments, output_file):
#     """
#     Create a PowerPoint presentation with bullet points distributed over slides.
#     """
#     # Initialize a presentation object
#     presentation = Presentation()

#     for i, bullets in enumerate(bullet_points_segments, 1):
#         # Add a slide for each segment
#         slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout

#         # Add title for the slide
#         title = slide.shapes.title
#         title.text = f"Segment {i}"

#         # Add bullet points to the slide
#         content = slide.placeholders[1]  # Content placeholder
#         for bullet in bullets:
#             p = content.text_frame.add_paragraph()
#             p.text = bullet
#             p.level = 0  # Top-level bullet point

#     # Save the presentation to the specified output file
#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")

# def main(file_path, threshold=0.5, output_file="output_presentation.pptx"):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation
#     create_presentation(bullet_points_segments, output_file)

# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "globel_warming.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path)
#**************************************************************************************************************************8

# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2


# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")


# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content


# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix


# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points


# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix


# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         # Split the segment into sentences or logical points
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments


# def create_presentation(bullet_points_segments, output_file, bullets_per_slide=5):
#     """
#     Create a PowerPoint presentation with bullet points distributed over slides.
#     """
#     # Initialize a presentation object
#     presentation = Presentation()

#     slide = None  # Current slide
#     current_bullet_count = 0  # Count of bullets on the current slide

#     for i, bullets in enumerate(bullet_points_segments, 1):
#         for bullet in bullets:
#             # If no slide or bullets exceed the limit, create a new slide
#             if not slide or current_bullet_count >= bullets_per_slide:
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
#                 slide.shapes.title.text = f"Segment {i}"  # Title for the new slide
#                 content = slide.placeholders[1]  # Content placeholder
#                 current_bullet_count = 0  # Reset bullet count for the new slide

#             # Add bullet to the current slide
#             p = content.text_frame.add_paragraph()
#             p.text = bullet
#             p.level = 0  # Top-level bullet point
#             current_bullet_count += 1

#     # Save the presentation to the specified output file
#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")


# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation with specified bullets per slide
#     create_presentation(bullet_points_segments, output_file, bullets_per_slide)


# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "globel_warming.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path, output_file="output_presentation.pptx", bullets_per_slide=5)

#*********************************************************************************************************************


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from transformers import pipeline


# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# # Step 2: Load Hugging Face's summarization pipeline for title generation
# summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content


# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix


# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points


# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix


# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         # Split the segment into sentences or logical points
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments


# def generate_title(segment):
#     """
#     Generate a title for a segment using Hugging Face's summarization pipeline.
#     """
#     summary = summarizer(segment, max_length=15, min_length=5, do_sample=False)
#     return summary[0]['summary_text'] if summary else "Untitled"


# def create_presentation(bullet_points_segments, segmented_text, output_file, bullets_per_slide=5):
#     """
#     Create a PowerPoint presentation with bullet points distributed over slides and titles generated from content.
#     """
#     # Initialize a presentation object
#     presentation = Presentation()

#     slide = None  # Current slide
#     current_bullet_count = 0  # Count of bullets on the current slide

#     for i, (bullets, segment) in enumerate(zip(bullet_points_segments, segmented_text), 1):
#         # Generate a title for the current segment using Hugging Face
#         title = generate_title(segment)
        
#         for bullet in bullets:
#             # If no slide or bullets exceed the limit, create a new slide
#             if not slide or current_bullet_count >= bullets_per_slide:
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
#                 slide.shapes.title.text = title  # Use the generated title as the slide title
#                 content = slide.placeholders[1]  # Content placeholder
#                 current_bullet_count = 0  # Reset bullet count for the new slide

#             # Add bullet to the current slide
#             p = content.text_frame.add_paragraph()
#             p.text = bullet
#             p.level = 0  # Top-level bullet point
#             current_bullet_count += 1

#     # Save the presentation to the specified output file
#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")


# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation with specified bullets per slide
#     create_presentation(bullet_points_segments, segmented_text, output_file, bullets_per_slide)


# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "globel_warming.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path, output_file="output_presentation.pptx", bullets_per_slide=5)


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from sklearn.feature_extraction.text import TfidfVectorizer
# import spacy

# # Step 1: Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")

# # Load spaCy for Named Entity Recognition
# nlp = spacy.load("en_core_web_sm")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content


# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix


# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         # Check for relative minima in similarity
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points


# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     # Step 2: Embed all paragraphs
#     embeddings = model.encode(paragraphs)

#     # Step 3: Compute cosine similarity
#     similarity_matrix = calculate_cosine_similarity(embeddings)

#     # Step 4: Find split points based on similarity minima
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     # Step 5: Segment the text
#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))  # Add the remaining text

#     return segmented_text, similarity_matrix


# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         # Split the segment into sentences or logical points
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments


# def generate_unique_title(segment):
#     """
#     Generate a unique, short title for the slide based on the content of the segment.
#     This uses keyword extraction via TF-IDF or Named Entity Recognition (NER).
#     """
#     # Use Named Entity Recognition (NER) from spaCy to identify key entities in the text
#     doc = nlp(segment)
#     entities = [ent.text for ent in doc.ents]  # Extract entities like names, places, etc.
    
#     # If entities are found, use the first entity as the title
#     if entities:
#         return entities[0]  # Take the first entity as a title
    
#     # If no entities are found, use TF-IDF to extract the top keywords
#     vectorizer = TfidfVectorizer(stop_words="english", max_features=1)  # Get the most important word
#     X = vectorizer.fit_transform([segment])
#     feature_names = vectorizer.get_feature_names_out()
#     return feature_names[0]  # Return the top keyword


# def create_presentation(bullet_points_segments, segmented_text, output_file, bullets_per_slide=5):
#     """
#     Create a PowerPoint presentation with bullet points distributed over slides and dynamic titles based on content.
#     """
#     # Initialize a presentation object
#     presentation = Presentation()

#     slide = None  # Current slide
#     current_bullet_count = 0  # Count of bullets on the current slide

#     for i, (bullets, segment) in enumerate(zip(bullet_points_segments, segmented_text), 1):
#         # Generate a unique title for each slide based on the segment content
#         title = generate_unique_title(segment)  # Use the generated title
        
#         for bullet in bullets:
#             # If no slide or bullets exceed the limit, create a new slide
#             if not slide or current_bullet_count >= bullets_per_slide:
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
#                 slide.shapes.title.text = title  # Use the unique title for the slide
#                 content = slide.placeholders[1]  # Content placeholder
#                 current_bullet_count = 0  # Reset bullet count for the new slide

#             # Add bullet to the current slide
#             p = content.text_frame.add_paragraph()
#             p.text = bullet
#             p.level = 0  # Top-level bullet point
#             current_bullet_count += 1

#     # Save the presentation to the specified output file
#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")


# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]  # Remove empty lines and extra spaces

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation with specified bullets per slide
#     create_presentation(bullet_points_segments, segmented_text, output_file, bullets_per_slide)


# # Example usage:
# if __name__ == "__main__":
#     # Specify the full path to your text or PDF file
#     file_path = "IRJET-V9I7590.pdf"  # Replace with your file path (e.g., "example.pdf" or "example.txt")
#     main(file_path, output_file="output_presentation.pptx", bullets_per_slide=5)


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from sklearn.feature_extraction.text import TfidfVectorizer
# import spacy
# from collections import Counter
# import re

# # Load models
# model = SentenceTransformer("all-mpnet-base-v2")
# nlp = spacy.load("en_core_web_sm")

# def read_file(file_path):
#     """Read content from text or PDF file."""
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     return content

# def segment_text(paragraphs, min_similarity=0.5):
#     """
#     Segment text into coherent sections using a simplified approach.
#     """
#     if not paragraphs:
#         return []
    
#     # Initialize segments with the first paragraph
#     segments = [paragraphs[0]]
#     current_segment = [paragraphs[0]]
    
#     # Process remaining paragraphs
#     for i in range(1, len(paragraphs)):
#         # Get embeddings for current paragraph and previous segment
#         current_embedding = model.encode([paragraphs[i]])[0]
#         prev_embedding = model.encode([" ".join(current_segment)])[0]
        
#         # Calculate similarity
#         similarity = np.dot(current_embedding, prev_embedding) / (
#             np.linalg.norm(current_embedding) * np.linalg.norm(prev_embedding)
#         )
        
#         # If similarity is high enough, add to current segment
#         if similarity >= min_similarity:
#             current_segment.append(paragraphs[i])
#         else:
#             # Start new segment
#             segments.append(" ".join(current_segment))
#             current_segment = [paragraphs[i]]
    
#     # Add the last segment
#     if current_segment:
#         segments.append(" ".join(current_segment))
    
#     return segments

# def split_into_sentences(text):
#     """Split text into sentences more accurately."""
#     # Clean the text first
#     text = re.sub(r'\s+', ' ', text).strip()
    
#     # Split on period followed by space and uppercase letter
#     sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
    
#     # Clean the sentences
#     sentences = [s.strip() for s in sentences if s.strip()]
#     return sentences

# def adjust_sentence_length(sentences, target=5):
#     """Adjust sentences to create exactly 5 bullets of similar length."""
#     if not sentences:
#         return []
    
#     total_text = ' '.join(sentences)
#     words = total_text.split()
    
#     # If we don't have enough words for 5 meaningful bullets
#     if len(words) < 10:  # Minimum 2 words per bullet
#         return sentences
    
#     # Calculate target words per bullet
#     words_per_bullet = max(2, len(words) // target)
    
#     adjusted_sentences = []
#     current_sentence = []
#     word_count = 0
    
#     for word in words:
#         current_sentence.append(word)
#         word_count += 1
        
#         if word_count >= words_per_bullet and len(adjusted_sentences) < target - 1:
#             adjusted_sentences.append(' '.join(current_sentence) + '.')
#             current_sentence = []
#             word_count = 0
    
#     # Add remaining words to the last bullet
#     if current_sentence:
#         adjusted_sentences.append(' '.join(current_sentence) + '.')
    
#     # If we have less than 5 bullets, duplicate some content
#     while len(adjusted_sentences) < target:
#         adjusted_sentences.append("Additional information: " + adjusted_sentences[0])
    
#     # If we have more than 5 bullets, combine extras with the last bullet
#     if len(adjusted_sentences) > target:
#         extra_content = ' '.join(adjusted_sentences[target-1:])
#         adjusted_sentences = adjusted_sentences[:target-1] + [extra_content]
    
#     return adjusted_sentences

# def generate_bullet_points(text):
#     """Convert text into exactly 5 bullet points."""
#     # Split into sentences
#     sentences = split_into_sentences(text)
    
#     # Adjust to exactly 5 sentences
#     adjusted_sentences = adjust_sentence_length(sentences)
    
#     # Convert to bullet points
#     bullets = [f"- {sentence}" for sentence in adjusted_sentences]
    
#     # Ensure exactly 5 bullets
#     while len(bullets) < 5:
#         bullets.append("- Additional point: " + bullets[0].replace("- ", ""))
    
#     return bullets[:5]  # Return exactly 5 bullets

# def generate_title(text):
#     """Generate a meaningful title from text content."""
#     doc = nlp(text)
    
#     # Try to get named entities first
#     entities = [ent.text for ent in doc.ents if len(ent.text.split()) <= 5]
#     if entities:
#         return entities[0][:50]
    
#     # If no entities, try noun chunks
#     chunks = [chunk.text for chunk in doc.noun_chunks if len(chunk.text.split()) <= 5]
#     if chunks:
#         return chunks[0][:50]
    
#     # Fallback to first sentence
#     first_sentence = text.split('.')[0]
#     return first_sentence[:50] + "..." if len(first_sentence) > 50 else first_sentence

# def create_presentation(text_segments, output_file):
#     """Create PowerPoint presentation with exactly 5 bullets per slide."""
#     prs = Presentation()
    
#     for segment in text_segments:
#         # Generate title and exactly 5 bullets
#         title = generate_title(segment)
#         bullets = generate_bullet_points(segment)
        
#         # Create slide
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text = title
        
#         # Add exactly 5 bullets to slide
#         tf = slide.shapes.placeholders[1].text_frame
#         for bullet in bullets:
#             p = tf.add_paragraph()
#             p.text = bullet
#             p.level = 0
    
#     prs.save(output_file)
#     print(f"Presentation saved as {output_file}")

# def main(file_path, output_file="output_presentation.pptx"):
#     """Main function to process document and create presentation."""
#     # Read content
#     content = read_file(file_path)
    
#     # Split into paragraphs
#     paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
    
#     # Segment the text
#     segments = segment_text(paragraphs)
    
#     # Create presentation
#     create_presentation(segments, output_file)

# if __name__ == "__main__":
#     file_path = "IRJET-V9I7590.pdf"
#     main(file_path, output_file="output_presentation.pptx")

#*******************************************************************************************

# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from sklearn.feature_extraction.text import TfidfVectorizer
# import spacy
# import re

# # Load the pre-trained embedding model
# model = SentenceTransformer("all-mpnet-base-v2")
# # Load spaCy model for NLP tasks
# nlp = spacy.load("en_core_web_sm")

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def generate_slide_title(bullets):
#     """
#     Generate an appropriate title for a slide based on its bullet points.
#     """
#     # Combine all bullet points into one text
#     combined_text = ' '.join([bullet.lstrip('- ') for bullet in bullets])
    
#     # Process the text with spaCy
#     doc = nlp(combined_text)
    
#     # Extract important phrases using TF-IDF
#     vectorizer = TfidfVectorizer(ngram_range=(1, 3), stop_words='english')
#     try:
#         tfidf_matrix = vectorizer.fit_transform([combined_text])
#         feature_names = vectorizer.get_feature_names_out()
#         scores = tfidf_matrix.toarray()[0]
        
#         # Get top phrases by TF-IDF score
#         top_indices = scores.argsort()[-3:][::-1]
#         important_phrases = [feature_names[i] for i in top_indices]
        
#         # Extract named entities
#         entities = [ent.text for ent in doc.ents 
#                    if ent.label_ in ['ORG', 'PRODUCT', 'GPE', 'PERSON', 'TOPIC', 'EVENT']]
        
#         # Combine entities and important phrases
#         title_candidates = []
#         if entities:
#             title_candidates.extend(entities[:1])
#         title_candidates.extend(important_phrases)
        
#         # Generate title
#         if title_candidates:
#             main_phrase = title_candidates[0].title()
#             if len(title_candidates) > 1 and len(main_phrase) < 30:
#                 return f"{main_phrase}: {title_candidates[1].title()}"
#             return main_phrase
#     except:
#         pass
    
#     # Fallback: use the first bullet point's main topic
#     first_bullet = bullets[0].lstrip('- ').split('.')[0]
#     if len(first_bullet) > 50:
#         return first_bullet[:47] + "..."
#     return first_bullet

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     embeddings = model.encode(paragraphs)
#     similarity_matrix = calculate_cosine_similarity(embeddings)
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))

#     return segmented_text, similarity_matrix

# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def create_presentation(bullet_points_segments, output_file, bullets_per_slide=5):
#     """
#     Create a PowerPoint presentation with generated titles and bullet points distributed over slides.
#     """
#     presentation = Presentation()
    
#     # Add title slide
#     title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
#     title = title_slide.shapes.title
#     subtitle = title_slide.placeholders[1]
#     title.text = "Document Overview"
#     subtitle.text = "Generated from Content Analysis"

#     current_segment_bullets = []
#     current_segment_index = 0

#     for i, segment_bullets in enumerate(bullet_points_segments, 1):
#         for bullet in segment_bullets:
#             current_segment_bullets.append(bullet)
            
#             # Create new slide when we reach the bullet limit or end of segment
#             if len(current_segment_bullets) >= bullets_per_slide or bullet == segment_bullets[-1]:
#                 # Generate title based on current bullet group
#                 slide_title = generate_slide_title(current_segment_bullets)
                
#                 # Create new slide
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                
#                 # Set the generated title
#                 title = slide.shapes.title
#                 title.text = slide_title
                
#                 # Add bullet points
#                 content = slide.placeholders[1]
#                 for bullet_point in current_segment_bullets:
#                     p = content.text_frame.add_paragraph()
#                     p.text = bullet_point.lstrip('- ')
#                     p.level = 0
                
#                 # Reset bullets for next slide
#                 current_segment_bullets = []
        
#         current_segment_index += 1

#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")

# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation with generated titles
#     create_presentation(bullet_points_segments, output_file, bullets_per_slide)

# if __name__ == "__main__":
#     file_path = "globel_warming.pdf"  # Replace with your file path
#     main(file_path, bullets_per_slide=5)

# $$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from transformers import pipeline, T5ForConditionalGeneration, T5Tokenizer
# import torch
# import re

# # Load the pre-trained embedding model for text segmentation
# model = SentenceTransformer("all-mpnet-base-v2")

# # Initialize the T5 model and tokenizer for title generation
# tokenizer = T5Tokenizer.from_pretrained('t5-base')
# title_model = T5ForConditionalGeneration.from_pretrained('t5-base')

# def read_file(file_path):
#     """
#     Read the content of a text or PDF file.
#     """
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    
#     return content

# def calculate_cosine_similarity(vectors):
#     """
#     Calculate cosine similarity for a list of vectors.
#     """
#     similarity_matrix = cosine_similarity(vectors)
#     return similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """
#     Identify split points based on the similarity matrix and a threshold.
#     """
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def generate_slide_title(bullets):
#     """
#     Generate an appropriate title for a slide using T5 model.
#     """
#     # Combine bullet points into a single text
#     text_content = ' '.join([bullet.lstrip('- ') for bullet in bullets])
    
#     # Prepare input for T5 model
#     input_text = f"summarize to title: {text_content}"
    
#     # Tokenize input
#     inputs = tokenizer.encode(input_text, return_tensors="pt", max_length=512, truncation=True)
    
#     # Generate title
#     outputs = title_model.generate(
#         inputs,
#         max_length=30,
#         min_length=5,
#         length_penalty=2.0,
#         num_beams=4,
#         early_stopping=True
#     )
    
#     # Decode and clean up the generated title
#     title = tokenizer.decode(outputs[0], skip_special_tokens=True)
    
#     # Clean and format the title
#     title = title.strip()
#     title = re.sub(r'\s+', ' ', title)  # Remove extra spaces
#     title = title.title()  # Capitalize words
    
#     # Handle empty or too short titles
#     if len(title) < 3:
#         return "Key Points"
    
#     return title

# def segment_text(paragraphs, threshold=0.5):
#     """
#     Segment text into meaningful paragraphs using embeddings and cosine similarity.
#     """
#     embeddings = model.encode(paragraphs)
#     similarity_matrix = calculate_cosine_similarity(embeddings)
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))

#     return segmented_text, similarity_matrix

# def generate_bullet_points(segmented_text):
#     """
#     Convert each segment into bullet points.
#     """
#     bullet_points_segments = []
#     for segment in segmented_text:
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def create_presentation(bullet_points_segments, output_file, bullets_per_slide=5):
#     """
#     Create a PowerPoint presentation with generated titles and bullet points distributed over slides.
#     """
#     presentation = Presentation()
    
#     # Add title slide
#     title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
#     title = title_slide.shapes.title
#     subtitle = title_slide.placeholders[1]
#     title.text = "Document Overview"
#     subtitle.text = "Generated from Content Analysis"

#     current_segment_bullets = []
#     current_segment_index = 0

#     for i, segment_bullets in enumerate(bullet_points_segments, 1):
#         for bullet in segment_bullets:
#             current_segment_bullets.append(bullet)
            
#             # Create new slide when we reach the bullet limit or end of segment
#             if len(current_segment_bullets) >= bullets_per_slide or bullet == segment_bullets[-1]:
#                 # Generate title based on current bullet group using T5
#                 slide_title = generate_slide_title(current_segment_bullets)
                
#                 # Create new slide
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                
#                 # Set the generated title
#                 title = slide.shapes.title
#                 title.text = slide_title
                
#                 # Add bullet points
#                 content = slide.placeholders[1]
#                 for bullet_point in current_segment_bullets:
#                     p = content.text_frame.add_paragraph()
#                     p.text = bullet_point.lstrip('- ')
#                     p.level = 0
                
#                 # Reset bullets for next slide
#                 current_segment_bullets = []
        
#         current_segment_index += 1

#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")

# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """
#     Main function to read file, segment text, generate bullet points, and create presentation.
#     """
#     # Read content from the file
#     content = read_file(file_path)
    
#     # Split the text into paragraphs
#     paragraphs = [para.strip() for para in content.split("\n") if para.strip()]

#     # Perform text segmentation
#     segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)

#     # Convert segments to bullet points
#     bullet_points_segments = generate_bullet_points(segmented_text)

#     # Create PowerPoint presentation with generated titles
#     create_presentation(bullet_points_segments, output_file, bullets_per_slide)

# if __name__ == "__main__":
#     file_path = "globel_warming.pdf"  # Replace with your file path
#     main(file_path, bullets_per_slide=5)


    #@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from pptx import Presentation
# import PyPDF2
# from transformers import pipeline, T5ForConditionalGeneration, T5Tokenizer
# from sklearn.feature_extraction.text import TfidfVectorizer
# import spacy
# import re
# import torch

# # Load required models
# model = SentenceTransformer("all-mpnet-base-v2")
# tokenizer = T5Tokenizer.from_pretrained('t5-base')
# title_model = T5ForConditionalGeneration.from_pretrained('t5-base')
# nlp = spacy.load("en_core_web_sm")

# def read_file(file_path):
#     """Read the content of a text or PDF file."""
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
#     return content

# def calculate_cosine_similarity(vectors):
#     """Calculate cosine similarity for a list of vectors."""
#     return cosine_similarity(vectors)

# def find_split_points(similarity_matrix, threshold=0.5):
#     """Identify split points based on the similarity matrix and a threshold."""
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def get_key_phrases(text, top_n=3):
#     """Extract key phrases using TF-IDF and named entities."""
#     # Process with spaCy
#     doc = nlp(text)
    
#     # Get named entities
#     entities = [ent.text for ent in doc.ents 
#                if ent.label_ in ['ORG', 'PRODUCT', 'GPE', 'PERSON', 'TOPIC', 'EVENT']]
    
#     # Extract phrases using TF-IDF
#     vectorizer = TfidfVectorizer(ngram_range=(1, 3), 
#                                 stop_words='english',
#                                 max_features=top_n)
#     try:
#         tfidf_matrix = vectorizer.fit_transform([text])
#         feature_names = vectorizer.get_feature_names_out()
#         scores = tfidf_matrix.toarray()[0]
#         top_indices = scores.argsort()[-top_n:][::-1]
#         phrases = [feature_names[i] for i in top_indices]
#     except:
#         phrases = []
    
#     return entities, phrases

# def generate_t5_title(text):
#     """Generate title using T5 model."""
#     input_text = f"summarize to brief title: {text}"
#     inputs = tokenizer.encode(input_text, return_tensors="pt", max_length=512, truncation=True)
    
#     outputs = title_model.generate(
#         inputs,
#         max_length=15,
#         min_length=3,
#         length_penalty=0.6,
#         num_beams=4,
#         early_stopping=True,
#         temperature=0.7
#     )
    
#     return tokenizer.decode(outputs[0], skip_special_tokens=True)

# def clean_title(title):
#     """Clean and format the title."""
#     # Remove extra spaces and newlines
#     title = re.sub(r'\s+', ' ', title).strip()
    
#     # Remove common prefixes
#     prefixes_to_remove = ['the ', 'a ', 'an ', 'this ', 'these ', 'those ']
#     for prefix in prefixes_to_remove:
#         if title.lower().startswith(prefix):
#             title = title[len(prefix):]
    
#     # Apply title case
#     title = title.title()
    
#     # Limit length
#     words = title.split()
#     if len(words) > 5:
#         title = ' '.join(words[:5])
    
#     return title

# def generate_slide_title(bullets, used_titles=None):
#     """Generate an appropriate title for a slide using multiple methods."""
#     if used_titles is None:
#         used_titles = set()
    
#     # Combine bullet points
#     text_content = ' '.join([bullet.lstrip('- ') for bullet in bullets])
    
#     # Get candidates from different methods
#     candidates = set()
    
#     # Method 1: T5 Model
#     t5_title = generate_t5_title(text_content)
#     if t5_title:
#         candidates.add(clean_title(t5_title))
    
#     # Method 2: Key phrases and entities
#     entities, phrases = get_key_phrases(text_content)
    
#     # Add entity-based titles
#     if entities:
#         for entity in entities[:2]:
#             candidates.add(clean_title(entity))
#             if phrases:
#                 candidates.add(clean_title(f"{entity}: {phrases[0]}"))
    
#     # Add phrase-based titles
#     if phrases:
#         candidates.add(clean_title(phrases[0]))
#         if len(phrases) > 1:
#             candidates.add(clean_title(f"{phrases[0]}: {phrases[1]}"))
    
#     # Remove any previously used titles
#     candidates = candidates - used_titles
    
#     # Select best candidate
#     if candidates:
#         title = min(candidates, key=len)  # Select the most concise title
#         used_titles.add(title)
#         return title
    
#     # Fallback: numbered generic title
#     base_title = "Key Points"
#     if base_title in used_titles:
#         counter = 1
#         while f"{base_title} {counter}" in used_titles:
#             counter += 1
#         title = f"{base_title} {counter}"
#     else:
#         title = base_title
    
#     used_titles.add(title)
#     return title

# def segment_text(paragraphs, threshold=0.5):
#     """Segment text into meaningful paragraphs using embeddings and cosine similarity."""
#     embeddings = model.encode(paragraphs)
#     similarity_matrix = calculate_cosine_similarity(embeddings)
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))

#     return segmented_text, similarity_matrix

# def generate_bullet_points(segmented_text):
#     """Convert each segment into bullet points."""
#     bullet_points_segments = []
#     for segment in segmented_text:
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def create_presentation(bullet_points_segments, output_file, bullets_per_slide=5):
#     """Create a PowerPoint presentation with generated titles and bullet points."""
#     presentation = Presentation()
#     used_titles = set()
    
#     # Add title slide
#     title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
#     title = title_slide.shapes.title
#     subtitle = title_slide.placeholders[1]
#     title.text = "Document Overview"
#     subtitle.text = "Generated from Content Analysis"

#     current_segment_bullets = []

#     for segment_bullets in bullet_points_segments:
#         for bullet in segment_bullets:
#             current_segment_bullets.append(bullet)
            
#             if len(current_segment_bullets) >= bullets_per_slide or bullet == segment_bullets[-1]:
#                 slide_title = generate_slide_title(current_segment_bullets, used_titles)
                
#                 slide = presentation.slides.add_slide(presentation.slide_layouts[1])
#                 title = slide.shapes.title
#                 title.text = slide_title
                
#                 content = slide.placeholders[1]
#                 for bullet_point in current_segment_bullets:
#                     p = content.text_frame.add_paragraph()
#                     p.text = bullet_point.lstrip('- ')
#                     p.level = 0
                
#                 current_segment_bullets = []

#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")

# def main(file_path, threshold=0.5, output_file="output_presentation.pptx", bullets_per_slide=5):
#     """Main function to process document and create presentation."""
#     try:
#         content = read_file(file_path)
#         paragraphs = [para.strip() for para in content.split("\n") if para.strip()]
#         segmented_text, similarity_matrix = segment_text(paragraphs, threshold=threshold)
#         bullet_points_segments = generate_bullet_points(segmented_text)
#         create_presentation(bullet_points_segments, output_file, bullets_per_slide)
#         print("Process completed successfully!")
#     except Exception as e:
#         print(f"An error occurred: {str(e)}")
#         raise

# if __name__ == "__main__":
#     file_path = "globel_warming.pdf"  # Replace with your file path
#     main(file_path, bullets_per_slide=5)


#$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$


# from sentence_transformers import SentenceTransformer
# import numpy as np
# from sklearn.metrics.pairwise import cosine_similarity
# from sklearn.cluster import AgglomerativeClustering
# from pptx import Presentation
# import PyPDF2
# from transformers import pipeline, T5ForConditionalGeneration, T5Tokenizer
# from sklearn.feature_extraction.text import TfidfVectorizer
# from pptx.util import Pt
# import spacy
# import re
# import torch

# # Load required models
# model = SentenceTransformer("all-mpnet-base-v2")
# tokenizer = T5Tokenizer.from_pretrained('t5-base')
# title_model = T5ForConditionalGeneration.from_pretrained('t5-base')
# nlp = spacy.load("en_core_web_sm")

# def read_file(file_path):
#     """Read the content of a text or PDF file."""
#     content = ""
#     if file_path.endswith(".txt"):
#         with open(file_path, "r", encoding="utf-8") as file:
#             content = file.read()
#     elif file_path.endswith(".pdf"):
#         with open(file_path, "rb") as file:
#             reader = PyPDF2.PdfReader(file)
#             for page in reader.pages:
#                 content += page.extract_text()
#     else:
#         raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
#     return content

# def group_similar_bullets(bullets, similarity_threshold=0.7):
#     """Group similar bullet points together using semantic similarity."""
#     if not bullets:
#         return []
    
#     # Get embeddings for all bullets
#     bullet_texts = [b.lstrip('- ') for b in bullets]
#     embeddings = model.encode(bullet_texts)
    
#     # Calculate similarity matrix
#     similarity_matrix = cosine_similarity(embeddings)
    
#     # Perform hierarchical clustering
#     n_clusters = max(1, len(bullets) // 5)  # Aim for about 5 bullets per group
#     clustering = AgglomerativeClustering(
#         n_clusters=n_clusters,
#         affinity='precomputed',
#         linkage='complete',
#         distance_threshold=None
#     )
    
#     # Convert similarity to distance
#     distance_matrix = 1 - similarity_matrix
#     cluster_labels = clustering.fit_predict(distance_matrix)
    
#     # Group bullets by cluster
#     grouped_bullets = {}
#     for i, label in enumerate(cluster_labels):
#         if label not in grouped_bullets:
#             grouped_bullets[label] = []
#         grouped_bullets[label].append(bullets[i])
    
#     return list(grouped_bullets.values())

# def get_topic_for_group(bullet_group):
#     """Extract the main topic for a group of bullet points."""
#     # Combine all bullets in the group
#     text = ' '.join([b.lstrip('- ') for b in bullet_group])
    
#     # Process with spaCy
#     doc = nlp(text)
    
#     # Extract key phrases using TF-IDF
#     vectorizer = TfidfVectorizer(ngram_range=(1, 2), 
#                                 stop_words='english',
#                                 max_features=3)
#     try:
#         tfidf_matrix = vectorizer.fit_transform([text])
#         feature_names = vectorizer.get_feature_names_out()
#         scores = tfidf_matrix.toarray()[0]
#         top_idx = scores.argmax()
#         main_topic = feature_names[top_idx].title()
#     except:
#         # Fallback to first sentence keywords
#         main_topic = ' '.join([token.text for token in list(doc)[:3]]).title()
    
#     return main_topic

# def clean_title(title):
#     """Clean and format the title."""
#     title = re.sub(r'\s+', ' ', title).strip()
#     title = title.title()
#     words = title.split()
#     if len(words) > 5:
#         title = ' '.join(words[:5])
#     return title

# def generate_slide_title(bullet_group, main_topic, used_titles):
#     """Generate an appropriate title for a slide based on the bullet group and topic."""
#     base_title = clean_title(main_topic)
    
#     if base_title not in used_titles:
#         used_titles.add(base_title)
#         return base_title
    
#     # Try with additional context
#     text_content = ' '.join([b.lstrip('- ') for b in bullet_group])
#     input_text = f"summarize to brief title: {text_content}"
#     inputs = tokenizer.encode(input_text, return_tensors="pt", max_length=512, truncation=True)
    
#     outputs = title_model.generate(
#         inputs,
#         max_length=15,
#         min_length=3,
#         length_penalty=0.6,
#         num_beams=4,
#         early_stopping=True,
#         temperature=0.7
#     )
    
#     alternative_title = clean_title(tokenizer.decode(outputs[0], skip_special_tokens=True))
    
#     if alternative_title not in used_titles:
#         used_titles.add(alternative_title)
#         return alternative_title
    
#     # Fallback to numbered version
#     counter = 1
#     while f"{base_title} {counter}" in used_titles:
#         counter += 1
#     final_title = f"{base_title} {counter}"
#     used_titles.add(final_title)
#     return final_title

# def segment_text(paragraphs, threshold=0.5):
#     """Segment text into meaningful paragraphs using embeddings and cosine similarity."""
#     embeddings = model.encode(paragraphs)
#     similarity_matrix = cosine_similarity(embeddings)
#     split_points = find_split_points(similarity_matrix, threshold=threshold)

#     segmented_text = []
#     start_idx = 0
#     for split_point in split_points:
#         segmented_text.append(" ".join(paragraphs[start_idx:split_point]))
#         start_idx = split_point
#     segmented_text.append(" ".join(paragraphs[start_idx:]))

#     return segmented_text, similarity_matrix

# def find_split_points(similarity_matrix, threshold=0.5):
#     """Identify split points based on the similarity matrix and a threshold."""
#     split_points = []
#     for i in range(1, len(similarity_matrix)):
#         if similarity_matrix[i-1][i-1] > similarity_matrix[i][i] and similarity_matrix[i][i] < threshold:
#             split_points.append(i)
#     return split_points

# def generate_bullet_points(segmented_text):
#     """Convert each segment into bullet points."""
#     bullet_points_segments = []
#     for segment in segmented_text:
#         sentences = segment.split(". ")
#         bullets = [f"- {sentence.strip()}." for sentence in sentences if sentence.strip()]
#         bullet_points_segments.append(bullets)
#     return bullet_points_segments

# def create_presentation(bullet_points_segments, output_file):
#     """Create a PowerPoint presentation with topically grouped bullets and appropriate titles."""
#     presentation = Presentation()
#     used_titles = set()
    
#     # Add title slide
#     title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
#     title = title_slide.shapes.title
#     subtitle = title_slide.placeholders[1]
#     title.text = "Document Overview"
#     subtitle.text = "Generated from Content Analysis"

#     # Flatten all bullets and group by similarity
#     all_bullets = [bullet for segment in bullet_points_segments for bullet in segment]
#     grouped_bullets = group_similar_bullets(all_bullets)
    
#     for bullet_group in grouped_bullets:
#         # Skip empty groups
#         if not bullet_group:
#             continue
            
#         # Get main topic for the group
#         main_topic = get_topic_for_group(bullet_group)
        
#         # Split into slides if too many bullets
#         for i in range(0, len(bullet_group), 5):
#             current_bullets = bullet_group[i:i+5]
            
#             # Generate slide title
#             slide_title = generate_slide_title(current_bullets, main_topic, used_titles)
            
#             # Create new slide
#             slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            
#             # Set title and topic heading
#             title = slide.shapes.title
#             title.text = slide_title
            
#             # Add bullet points
#             content = slide.placeholders[1]
            
#             # Add topic as heading if it's the first slide for this group
#             if i == 0:
#                 p = content.text_frame.add_paragraph()
#                 p.text = f"Topic: {main_topic}"
#                 p.font.bold = True
#                 p.font.size = Pt(14)
            
#             # Add bullets
#             for bullet in current_bullets:
#                 p = content.text_frame.add_paragraph()
#                 p.text = bullet.lstrip('- ')
#                 p.level = 0

#     presentation.save(output_file)
#     print(f"Presentation created successfully: {output_file}")

# def main(file_path, output_file="output_presentation.pptx"):
#     """Main function to process document and create presentation."""
#     try:
#         print("Reading file...")
#         content = read_file(file_path)
        
#         print("Processing text...")
#         paragraphs = [para.strip() for para in content.split("\n") if para.strip()]
        
#         print("Segmenting text...")
#         segmented_text, similarity_matrix = segment_text(paragraphs)
        
#         print("Generating bullet points...")
#         bullet_points_segments = generate_bullet_points(segmented_text)
        
#         print("Creating presentation...")
#         create_presentation(bullet_points_segments, output_file)
        
#         print("Process completed successfully!")
#     except Exception as e:
#         print(f"An error occurred: {str(e)}")
#         raise

# if __name__ == "__main__":
#     file_path = "pollution.pdf"  # Replace with your file path
#     main(file_path)


#$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$


from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.util import Pt
import PyPDF2
from transformers import pipeline
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import NMF
import spacy
import re

# Load required models
model = SentenceTransformer("all-mpnet-base-v2")
nlp = spacy.load("en_core_web_sm")
zero_shot_classifier = pipeline("zero-shot-classification")

def read_file(file_path):
    """Read the content of a text or PDF file."""
    content = ""
    if file_path.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()
    elif file_path.endswith(".pdf"):
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                content += page.extract_text()
    else:
        raise ValueError("Unsupported file format. Please provide a .txt or .pdf file.")
    return content

def extract_topics(text, num_topics=10):
    """Extract main topics from the text using NMF."""
    # Preprocess text
    vectorizer = TfidfVectorizer(
        max_features=1000,
        stop_words='english',
        ngram_range=(1, 2)
    )
    
    # Create document-term matrix
    dtm = vectorizer.fit_transform([text])
    
    # Apply NMF
    nmf_model = NMF(n_components=num_topics, random_state=42)
    nmf_output = nmf_model.fit_transform(dtm)
    
    # Get feature names (terms)
    feature_names = vectorizer.get_feature_names_out()
    
    # Extract top terms for each topic
    topics = []
    for topic_idx in range(num_topics):
        top_terms_idx = nmf_model.components_[topic_idx].argsort()[:-6:-1]
        top_terms = [feature_names[i] for i in top_terms_idx]
        topic_name = " ".join(top_terms[:2])  # Use first two terms as topic name
        topics.append(topic_name)
    
    return topics

def split_into_sentences(text):
    """Split text into sentences using spaCy."""
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
    return [s for s in sentences if len(s.split()) >= 4]  # Filter out very short sentences

def get_sentence_embeddings(sentences):
    """Get embeddings for sentences."""
    return model.encode(sentences)

def classify_sentence_topic(sentence, topics, sentence_embedding, topic_embeddings):
    """Classify which topic the sentence belongs to using semantic similarity."""
    sentence_similarity = cosine_similarity([sentence_embedding], topic_embeddings)[0]
    best_topic_idx = np.argmax(sentence_similarity)
    confidence = sentence_similarity[best_topic_idx]
    return topics[best_topic_idx], confidence

def group_sentences_by_topic(sentences, topics):
    """Group sentences by detected topics using semantic similarity."""
    # Get embeddings
    sentence_embeddings = get_sentence_embeddings(sentences)
    topic_embeddings = get_sentence_embeddings(topics)
    
    topic_groups = {topic: [] for topic in topics}
    
    for sentence, sentence_embedding in zip(sentences, sentence_embeddings):
        topic, confidence = classify_sentence_topic(
            sentence, 
            topics, 
            sentence_embedding, 
            topic_embeddings
        )
        
        if confidence > 0.3:  # Lower threshold for broader topic matching
            topic_groups[topic].append(sentence)
    
    return topic_groups

def split_content_for_slides(sentences, max_bullets_per_slide=5):
    """Split content into chunks appropriate for slides."""
    slides_content = []
    for i in range(0, len(sentences), max_bullets_per_slide):
        slides_content.append(sentences[i:i + max_bullets_per_slide])
    return slides_content

def clean_topic_name(topic):
    """Clean and format topic name for slide title."""
    # Remove special characters and extra spaces
    cleaned = re.sub(r'[^\w\s]', '', topic)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    # Capitalize first letter of each word
    cleaned = cleaned.title()
    
    # Limit length
    words = cleaned.split()
    if len(words) > 4:
        cleaned = ' '.join(words[:4])
    
    return cleaned

def detect_subtopic_type(sentences):
    """Detect if the content group is about concepts, examples, or applications."""
    combined_text = ' '.join(sentences).lower()
    
    # Define keyword patterns for different subtypes
    patterns = {
        'Definition & Concepts': ['define', 'concept', 'refer', 'mean', 'is a', 'are the'],
        'Examples & Cases': ['example', 'instance', 'case', 'such as', 'like'],
        'Applications & Impact': ['apply', 'impact', 'effect', 'result', 'cause', 'lead to'],
        'Analysis & Details': ['analyze', 'examine', 'consist', 'contain', 'include']
    }
    
    # Count pattern matches
    matches = {subtype: sum(1 for pattern in patterns[subtype] if pattern in combined_text)
              for subtype in patterns}
    
    # Return the subtype with most matches, default to Analysis & Details
    return max(matches.items(), key=lambda x: x[1])[0]

def create_presentation(topic_groups, output_file):
    """Create a PowerPoint presentation with automatically detected topics."""
    presentation = Presentation()
    
    # Add title slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Topic Analysis"
    subtitle.text = "Key Concepts and Insights"
    
    # Create slides for each topic
    for topic, sentences in topic_groups.items():
        if not sentences:  # Skip empty topics
            continue
        
        # Split content into multiple slides if needed
        slide_contents = split_content_for_slides(sentences)
        
        for idx, content in enumerate(slide_contents):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            
            # Set title
            title = slide.shapes.title
            clean_title = clean_topic_name(topic)
            title_text = clean_title
            if len(slide_contents) > 5:
                title_text += f" (Part {idx + 1})"
            title.text = title_text
            
            # Add content
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            
            # Add subheading for first slide of each topic
            if idx == 0:
                p = text_frame.add_paragraph()
                p.text = detect_subtopic_type(content)
                p.font.bold = True
                p.font.size = Pt(14)
            
            # Add bullet points
            for sentence in content:
                p = text_frame.add_paragraph()
                p.text = sentence
                p.level = 0
    
    presentation.save(output_file)
    print(f"Presentation created successfully: {output_file}")

def main(file_path, output_file="auto_generated_presentation.pptx", num_topics=5):
    """Main function to process document and create presentation."""
    try:
        print("Reading file...")
        content = read_file(file_path)
        
        print("Detecting main topics...")
        topics = extract_topics(content, num_topics)
        print(f"Detected topics: {topics}")
        
        print("Processing text...")
        sentences = split_into_sentences(content)
        
        print("Classifying content by topics...")
        topic_groups = group_sentences_by_topic(sentences, topics)
        
        print("Creating presentation...")
        create_presentation(topic_groups, output_file)
        
        print("Process completed successfully!")
    except Exception as e:
        print(f"An error occurred: {str(e)}")
        raise

if __name__ == "__main__":
    file_path = "pollution.pdf"  # Replace with your file path
    main(file_path)


